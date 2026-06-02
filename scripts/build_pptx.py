from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a 16:9 PPTX from rendered slide images.")
    parser.add_argument("--manifest", required=True, help="Path to slide-manifest.json.")
    parser.add_argument("--out", required=True, help="Output .pptx path.")
    parser.add_argument("--video", action="append", default=[], help="Optional video path to embed on the last slide.")
    parser.add_argument("--width", type=float, default=13.333333, help="Slide width in inches.")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height in inches.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    manifest_path = Path(args.manifest).resolve()
    output_path = Path(args.out).resolve()
    manifest_dir = manifest_path.parent
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(args.width)
    prs.slide_height = Inches(args.height)
    blank_layout = prs.slide_layouts[6]

    videos = [Path(item).resolve() for item in args.video]
    last_video = videos[0] if videos else None

    for item in manifest["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        image_path = (manifest_dir / item["image"]).resolve()
        is_last = item["index"] == len(manifest["slides"])

        if is_last and last_video and last_video.exists():
            try:
                slide.shapes.add_movie(
                    str(last_video),
                    0,
                    0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                    poster_frame_image=str(image_path),
                    mime_type="video/mp4",
                )
                continue
            except Exception:
                pass

        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(output_path))
    print(f"Saved PPTX: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

